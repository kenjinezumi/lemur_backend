from pptx import Presentation
from pptx.util import Inches, Pt
import io

def create_powerpoint(content):
    """
    Create a PowerPoint presentation from the given content.

    Args:
        content (dict): The content to add to the slides.

    Returns:
        bytes: The generated PowerPoint presentation as a byte array.
    """
    prs = Presentation()

    for slide_no, slide_content in content.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        add_content_to_slide(slide, slide_no, slide_content)

    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    return pptx_file.getvalue()

def add_content_to_slide(slide, slide_no, content):
    """
    Add the given content to the specified slide.

    Args:
        slide (Slide): The slide to add content to.
        slide_no (int): The slide number.
        content (dict): The content to add to the slide.
    """
    title = slide.shapes.title
    title.text = f"Slide {slide_no}"

    y_offset = 1.5
    if content.get('insights'):
        add_textbox(slide, 'Insights', content['insights'], y_offset)
        y_offset += len(content['insights']) * 0.5 + 0.5

    if content.get('recommendations'):
        add_textbox(slide, 'Recommendations', content['recommendations'], y_offset)
        y_offset += len(content['recommendations']) * 0.5 + 0.5

    if content.get('drivers'):
        add_textbox(slide, 'Drivers', content['drivers'], y_offset)
        y_offset += len(content['drivers']) * 0.5 + 0.5

    if content.get('codes'):
        add_textbox(slide, 'Codes', content['codes'], y_offset)

def add_textbox(slide, header, items, y_offset):
    """
    Add a textbox with the given header and items to the slide at the specified vertical position.

    Args:
        slide (Slide): The slide to add the textbox to.
        header (str): The header text.
        items (list): The list of items to add to the textbox.
        y_offset (float): The vertical position to place the textbox.
    """
    left = Inches(0.5)
    top = Inches(y_offset)
    width = Inches(9)
    height = Inches(0.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = f"{header}:"
    p.font.bold = True
    p.font.size = Pt(18)

    for item in items:
        p = text_frame.add_paragraph()
        p.text = f"- {item}"
        p.font.size = Pt(14)
