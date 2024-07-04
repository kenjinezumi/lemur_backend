import os
import re
import logging
import json
import requests
from flask import Flask, request, jsonify
from google.auth import default
from googleapiclient.discovery import build
import time
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from googleapiclient.http import MediaFileUpload
from flasgger import Swagger

# from google.oauth2 import service_account

# Flask app
app = Flask(__name__)
swagger = Swagger(app)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Google Drive API setup
SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/presentations",
]
# creds, project = default(scopes=SCOPES)
# SERVICE_ACCOUNT_FILE = '/app/service-account.json'  # Update with your file path
# creds = service_account.Credentials.from_service_account_file(
#     SERVICE_ACCOUNT_FILE, scopes=SCOPES)
creds, project = default(scopes=SCOPES)
drive_service = build("drive", "v3", credentials=creds)

def fetch_slide_data_with_retry(api_url, slide_no, retries=3):
    attempt = 0
    while attempt < retries:
        try:
            response = requests.post(
                api_url, json={"slide_no": str(slide_no)}, timeout=3600
            )
            logger.info(
                f"API response status code for slide {slide_no}: {response.status_code}"
            )
            logger.info(f"API response content for slide {slide_no}: {response.text}")
            response.raise_for_status()  # Raise an exception for HTTP errors
            return response.json()
        except requests.exceptions.RequestException as e:
            logger.error(f"Attempt {attempt + 1} failed with error: {e}")
            attempt += 1
            time.sleep(2**attempt)  # Exponential backoff
    raise Exception(
        f"Failed to fetch slide data for slide {slide_no} after several retries"
    )

# Hardcoded data for slides
hardcoded_data = {
    14: {"data": {"GLOBAL": {"Direct Named": {"QSO": {"QTD": "1.18K", "Attain": "0.12%", "YoY": "1.65%"}, "Pipeline": {"QTD": "$3077.79M", "Attain": "2.6%", "YoY": "-2.03%"}}, "SMB": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "0.88%"}, "Pipeline": {"QTD": "$1033.29M", "Attain": "2.64%", "YoY": "-1.22%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-0.72%"}, "Pipeline": {"QTD": "$1043.98M", "Attain": "2.61%", "YoY": "-0.73%"}}, "Partner": {"Pipeline": {"QTD": "$6201.1M", "Attain": "2.57%", "YoY": "-0.47%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.96K", "Attain": "0.12%", "YoY": "1.18%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12370.18M", "Attain": "2.59%", "YoY": "-1.4%"}}}, "NORTHAM": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "0.06%"}, "Pipeline": {"QTD": "$3123.27M", "Attain": "2.59%", "YoY": "2.9%"}}, "SMB": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-1.62%"}, "Pipeline": {"QTD": "$1039.55M", "Attain": "2.57%", "YoY": "3.94%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-2.73%"}, "Pipeline": {"QTD": "$1031.71M", "Attain": "2.56%", "YoY": "-1.33%"}}, "Partner": {"Pipeline": {"QTD": "$6171.58M", "Attain": "2.58%", "YoY": "-0.87%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "-0.83%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12379.28M", "Attain": "2.59%", "YoY": "0.14%"}}}, "US PUBLIC SECTOR": {"Direct Named": {"QSO": {"QTD": "1.16K", "Attain": "0.12%", "YoY": "-1.98%"}, "Pipeline": {"QTD": "$3118.48M", "Attain": "2.62%", "YoY": "-0.63%"}}, "Partner": {"Pipeline": {"QTD": "$6303.99M", "Attain": "2.59%", "YoY": "1.33%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "-1.99%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12534.91M", "Attain": "2.6%", "YoY": "0.17%"}}}, "JAPAC": {"Direct Named": {"QSO": {"QTD": "3.48K", "Attain": "0.12%", "YoY": "-0.71%"}, "Pipeline": {"QTD": "$9338.47M", "Attain": "2.59%", "YoY": "-0.66%"}}, "SMB": {"QSO": {"QTD": "1.15K", "Attain": "0.11%", "YoY": "-1.28%"}, "Pipeline": {"QTD": "$3073.57M", "Attain": "2.6%", "YoY": "-1.26%"}}, "Startup": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "1.3%"}, "Pipeline": {"QTD": "$3135.85M", "Attain": "2.58%", "YoY": "1.89%"}}, "Partner": {"Pipeline": {"QTD": "$18650.57M", "Attain": "2.59%", "YoY": "0.04%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "5.81K", "Attain": "0.12%", "YoY": "-0.57%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$37300.23M", "Attain": "2.59%", "YoY": "-0.29%"}}}, "EMEA": {"Direct Named": {"QSO": {"QTD": "1.16K", "Attain": "0.12%", "YoY": "-0.63%"}, "Pipeline": {"QTD": "$3104.1M", "Attain": "2.58%", "YoY": "0.34%"}}, "SMB": {"QSO": {"QTD": "0.4K", "Attain": "0.12%", "YoY": "2.95%"}, "Pipeline": {"QTD": "$1069.47M", "Attain": "2.56%", "YoY": "3.32%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "0.99%"}, "Pipeline": {"QTD": "$1011.2M", "Attain": "2.54%", "YoY": "-3.36%"}}, "Partner": {"Pipeline": {"QTD": "$6228.18M", "Attain": "2.57%", "YoY": "1.36%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.95K", "Attain": "0.12%", "YoY": "-0.2%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12444.09M", "Attain": "2.58%", "YoY": "0.58%"}}}, "LATAM": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "-0.57%"}, "Pipeline": {"QTD": "$3067.74M", "Attain": "2.57%", "YoY": "-2.15%"}}, "SMB": {"QSO": {"QTD": "0.4K", "Attain": "0.12%", "YoY": "3.07%"}, "Pipeline": {"QTD": "$1044.78M", "Attain": "2.6%", "YoY": "0.51%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-2.03%"}, "Pipeline": {"QTD": "$1046.88M", "Attain": "2.62%", "YoY": "1.28%"}}, "Partner": {"Pipeline": {"QTD": "$6250.27M", "Attain": "2.6%", "YoY": "-0.44%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.96K", "Attain": "0.12%", "YoY": "-0.33%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12451.56M", "Attain": "2.6%", "YoY": "-0.16%"}}}}, "insights": [{"title": "EMEA SMB Direct Named Program Underperforming in QTD QSO Pacing", "narrative": "The *EMEA SMB Direct Named* marketing program is significantly underperforming, achieving a QTD QSO pacing of only 0.17% against a target of 95%. This indicates a critical bottleneck in converting inquiries from key channels like *Display - Paid Social* (12.59% of inquiries) and *Email* (13.04% of inquiries) to QSOs.  A thorough analysis of campaign-level conversion rates within this program, particularly focusing on sales follow-up rates, is crucial to identify the root causes and implement corrective actions."}, {"title": "EMEA SMB Partner QTD QSO Pacing Lags Despite Slight YoY Improvement", "narrative": "EMEA SMB Partner QTD QSO Pacing is alarmingly low at **0.17%**, signaling potential difficulties in achieving quarterly targets despite a marginal **0.67%** YoY increase.  This underperformance is further emphasized by the substantial *$11.3M* pipeline generated by campaigns like *P&C Top Summit January 2024*, which unfortunately struggles to translate into qualified opportunities due to a low QSO conversion rate. To address this, prioritize optimizing pipeline conversion by analyzing high-performing campaigns like *'24 Gartner Supply Chain Symposium/Xpo'* (**50.63%** SAL Conversion Rate) and replicating their successful strategies within the EMEA SMB Partner segment. Additionally, benchmarking the performance of EMEA SMB Partner marketing programs against successful initiatives in other regions like NORTHAM or PUBLIC SECTOR, such as *Cloud Architecture Framework: Made in The Cloud*, can provide valuable insights for improvement."}]},
    15: {"data": {"GLOBAL": {"Direct Named": {"QSO": {"QTD": "1.18K", "Attain": "0.12%", "YoY": "1.65%"}, "Pipeline": {"QTD": "$3077.79M", "Attain": "2.6%", "YoY": "-2.03%"}}, "SMB": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "0.88%"}, "Pipeline": {"QTD": "$1033.29M", "Attain": "2.64%", "YoY": "-1.22%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-0.72%"}, "Pipeline": {"QTD": "$1043.98M", "Attain": "2.61%", "YoY": "-0.73%"}}, "Partner": {"Pipeline": {"QTD": "$6201.1M", "Attain": "2.57%", "YoY": "-0.47%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.96K", "Attain": "0.12%", "YoY": "1.18%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12370.18M", "Attain": "2.59%", "YoY": "-1.4%"}}}, "NORTHAM": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "0.06%"}, "Pipeline": {"QTD": "$3123.27M", "Attain": "2.59%", "YoY": "2.9%"}}, "SMB": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-1.62%"}, "Pipeline": {"QTD": "$1039.55M", "Attain": "2.57%", "YoY": "3.94%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-2.73%"}, "Pipeline": {"QTD": "$1031.71M", "Attain": "2.56%", "YoY": "-1.33%"}}, "Partner": {"Pipeline": {"QTD": "$6171.58M", "Attain": "2.58%", "YoY": "-0.87%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "-0.83%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12379.28M", "Attain": "2.59%", "YoY": "0.14%"}}}, "US PUBLIC SECTOR": {"Direct Named": {"QSO": {"QTD": "1.16K", "Attain": "0.12%", "YoY": "-1.98%"}, "Pipeline": {"QTD": "$3118.48M", "Attain": "2.62%", "YoY": "-0.63%"}}, "Partner": {"Pipeline": {"QTD": "$6303.99M", "Attain": "2.59%", "YoY": "1.33%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "-1.99%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12534.91M", "Attain": "2.6%", "YoY": "0.17%"}}}, "JAPAC": {"Direct Named": {"QSO": {"QTD": "3.48K", "Attain": "0.12%", "YoY": "-0.71%"}, "Pipeline": {"QTD": "$9338.47M", "Attain": "2.59%", "YoY": "-0.66%"}}, "SMB": {"QSO": {"QTD": "1.15K", "Attain": "0.11%", "YoY": "-1.28%"}, "Pipeline": {"QTD": "$3073.57M", "Attain": "2.6%", "YoY": "-1.26%"}}, "Startup": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "1.3%"}, "Pipeline": {"QTD": "$3135.85M", "Attain": "2.58%", "YoY": "1.89%"}}, "Partner": {"Pipeline": {"QTD": "$18650.57M", "Attain": "2.59%", "YoY": "0.04%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "5.81K", "Attain": "0.12%", "YoY": "-0.57%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$37300.23M", "Attain": "2.59%", "YoY": "-0.29%"}}}, "EMEA": {"Direct Named": {"QSO": {"QTD": "1.16K", "Attain": "0.12%", "YoY": "-0.63%"}, "Pipeline": {"QTD": "$3104.1M", "Attain": "2.58%", "YoY": "0.34%"}}, "SMB": {"QSO": {"QTD": "0.4K", "Attain": "0.12%", "YoY": "2.95%"}, "Pipeline": {"QTD": "$1069.47M", "Attain": "2.56%", "YoY": "3.32%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "0.99%"}, "Pipeline": {"QTD": "$1011.2M", "Attain": "2.54%", "YoY": "-3.36%"}}, "Partner": {"Pipeline": {"QTD": "$6228.18M", "Attain": "2.57%", "YoY": "1.36%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.95K", "Attain": "0.12%", "YoY": "-0.2%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12444.09M", "Attain": "2.58%", "YoY": "0.58%"}}}, "LATAM": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "-0.57%"}, "Pipeline": {"QTD": "$3067.74M", "Attain": "2.57%", "YoY": "-2.15%"}}, "SMB": {"QSO": {"QTD": "0.4K", "Attain": "0.12%", "YoY": "3.07%"}, "Pipeline": {"QTD": "$1044.78M", "Attain": "2.6%", "YoY": "0.51%"}}, "Startup": {"QSO": {"QTD": "0.39K", "Attain": "0.12%", "YoY": "-2.03%"}, "Pipeline": {"QTD": "$1046.88M", "Attain": "2.62%", "YoY": "1.28%"}}, "Partner": {"Pipeline": {"QTD": "$6250.27M", "Attain": "2.6%", "YoY": "-0.44%"}}, "GCP Direct QSOs": {"QSO": {"QTD": "1.96K", "Attain": "0.12%", "YoY": "-0.33%"}}, "GCP Direct + Partner Pipe": {"Pipeline": {"QTD": "$12451.56M", "Attain": "2.6%", "YoY": "-0.16%"}}}}, "insights": [{"title": "EMEA SMB Direct Named Program Underperforming in QTD QSO Pacing", "narrative": "The *EMEA SMB Direct Named* marketing program is significantly underperforming, achieving a QTD QSO pacing of only 0.17% against a target of 95%. This indicates a critical bottleneck in converting inquiries from key channels like *Display - Paid Social* (12.59% of inquiries) and *Email* (13.04% of inquiries) to QSOs.  A thorough analysis of campaign-level conversion rates within this program, particularly focusing on sales follow-up rates, is crucial to identify the root causes and implement corrective actions."}, {"title": "EMEA SMB Partner QTD QSO Pacing Lags Despite Slight YoY Improvement", "narrative": "EMEA SMB Partner QTD QSO Pacing is alarmingly low at **0.17%**, signaling potential difficulties in achieving quarterly targets despite a marginal **0.67%** YoY increase.  This underperformance is further emphasized by the substantial *$11.3M* pipeline generated by campaigns like *P&C Top Summit January 2024*, which unfortunately struggles to translate into qualified opportunities due to a low QSO conversion rate. To address this, prioritize optimizing pipeline conversion by analyzing high-performing campaigns like *'24 Gartner Supply Chain Symposium/Xpo'* (**50.63%** SAL Conversion Rate) and replicating their successful strategies within the EMEA SMB Partner segment. Additionally, benchmarking the performance of EMEA SMB Partner marketing programs against successful initiatives in other regions like NORTHAM or PUBLIC SECTOR, such as *Cloud Architecture Framework: Made in The Cloud*, can provide valuable insights for improvement."}]},
    16: {"data": {"GLOBAL": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "1.23%"}, "Pipeline": {"QTD": "$3117.27M", "Attain": "2.63%", "YoY": "0.15%"}}, "Partner": {"Pipeline": {"QTD": "$6179.85M", "Attain": "2.59%", "YoY": "-0.43%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "0.26%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12444.05M", "Attain": "2.61%", "YoY": "0.17%"}}}, "NORTHAM": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "0.27%"}, "Pipeline": {"QTD": "$3093.7M", "Attain": "2.59%", "YoY": "0.71%"}}, "Partner": {"Pipeline": {"QTD": "$6225.05M", "Attain": "2.59%", "YoY": "0.27%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "0.33%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12415.4M", "Attain": "2.59%", "YoY": "0.29%"}}}, "US PUBLIC SECTOR": {"Direct Named": {"QSO": {"QTD": "1.19K", "Attain": "0.12%", "YoY": "2.56%"}, "Pipeline": {"QTD": "$3115.76M", "Attain": "2.58%", "YoY": "1.04%"}}, "Partner": {"Pipeline": {"QTD": "$6227.15M", "Attain": "2.59%", "YoY": "0.62%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.95K", "Attain": "0.12%", "YoY": "0.09%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12465.17M", "Attain": "2.59%", "YoY": "0.6%"}}}, "JAPAC": {"Direct Named": {"QSO": {"QTD": "3.49K", "Attain": "0.12%", "YoY": "-1.22%"}, "Pipeline": {"QTD": "$9435.17M", "Attain": "2.62%", "YoY": "1.41%"}}, "Partner": {"Pipeline": {"QTD": "$18748.45M", "Attain": "2.61%", "YoY": "0.24%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "5.84K", "Attain": "0.12%", "YoY": "-0.38%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$37564.69M", "Attain": "2.61%", "YoY": "0.63%"}}}, "EMEA": {"Direct Named": {"QSO": {"QTD": "1.2K", "Attain": "0.12%", "YoY": "1.28%"}, "Pipeline": {"QTD": "$3118.31M", "Attain": "2.61%", "YoY": "-0.74%"}}, "Partner": {"Pipeline": {"QTD": "$6256.96M", "Attain": "2.59%", "YoY": "0.25%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.97K", "Attain": "0.12%", "YoY": "1.1%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12485.45M", "Attain": "2.6%", "YoY": "0.02%"}}}, "LATAM": {"Direct Named": {"QSO": {"QTD": "1.16K", "Attain": "0.11%", "YoY": "-1.57%"}, "Pipeline": {"QTD": "$3118.28M", "Attain": "2.6%", "YoY": "0.05%"}}, "Partner": {"Pipeline": {"QTD": "$6152.72M", "Attain": "2.57%", "YoY": "-0.71%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.95K", "Attain": "0.12%", "YoY": "-1.16%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12434.49M", "Attain": "2.6%", "YoY": "-0.48%"}}}}, "insights": [{"title": "EMEA SMB Direct Named Program Underperforming in QTD QSO Pacing", "narrative": "The *EMEA SMB Direct Named* marketing program is significantly underperforming, achieving a QTD QSO pacing of only 0.17% against a target of 95%. This indicates a critical bottleneck in converting inquiries from key channels like *Display - Paid Social* (12.59% of inquiries) and *Email* (13.04% of inquiries) to QSOs.  A thorough analysis of campaign-level conversion rates within this program, particularly focusing on sales follow-up rates, is crucial to identify the root causes and implement corrective actions."}, {"title": "EMEA SMB Partner QTD QSO Pacing Lags Despite Slight YoY Improvement", "narrative": "EMEA SMB Partner QTD QSO Pacing is alarmingly low at **0.17%**, signaling potential difficulties in achieving quarterly targets despite a marginal **0.67%** YoY increase.  This underperformance is further emphasized by the substantial *$11.3M* pipeline generated by campaigns like *P&C Top Summit January 2024*, which unfortunately struggles to translate into qualified opportunities due to a low QSO conversion rate. To address this, prioritize optimizing pipeline conversion by analyzing high-performing campaigns like *'24 Gartner Supply Chain Symposium/Xpo'* (**50.63%** SAL Conversion Rate) and replicating their successful strategies within the EMEA SMB Partner segment. Additionally, benchmarking the performance of EMEA SMB Partner marketing programs against successful initiatives in other regions like NORTHAM or PUBLIC SECTOR, such as *Cloud Architecture Framework: Made in The Cloud*, can provide valuable insights for improvement."}]},
    17: {"data": {"GLOBAL": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "1.23%"}, "Pipeline": {"QTD": "$3117.27M", "Attain": "2.63%", "YoY": "0.15%"}}, "Partner": {"Pipeline": {"QTD": "$6179.85M", "Attain": "2.59%", "YoY": "-0.43%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "0.26%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12444.05M", "Attain": "2.61%", "YoY": "0.17%"}}}, "NORTHAM": {"Direct Named": {"QSO": {"QTD": "1.17K", "Attain": "0.12%", "YoY": "0.27%"}, "Pipeline": {"QTD": "$3093.7M", "Attain": "2.59%", "YoY": "0.71%"}}, "Partner": {"Pipeline": {"QTD": "$6225.05M", "Attain": "2.59%", "YoY": "0.27%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.94K", "Attain": "0.12%", "YoY": "0.33%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12415.4M", "Attain": "2.59%", "YoY": "0.29%"}}}, "US PUBLIC SECTOR": {"Direct Named": {"QSO": {"QTD": "1.19K", "Attain": "0.12%", "YoY": "2.56%"}, "Pipeline": {"QTD": "$3115.76M", "Attain": "2.58%", "YoY": "1.04%"}}, "Partner": {"Pipeline": {"QTD": "$6227.15M", "Attain": "2.59%", "YoY": "0.62%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.95K", "Attain": "0.12%", "YoY": "0.09%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12465.17M", "Attain": "2.59%", "YoY": "0.6%"}}}, "JAPAC": {"Direct Named": {"QSO": {"QTD": "3.49K", "Attain": "0.12%", "YoY": "-1.22%"}, "Pipeline": {"QTD": "$9435.17M", "Attain": "2.62%", "YoY": "1.41%"}}, "Partner": {"Pipeline": {"QTD": "$18748.45M", "Attain": "2.61%", "YoY": "0.24%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "5.84K", "Attain": "0.12%", "YoY": "-0.38%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$37564.69M", "Attain": "2.61%", "YoY": "0.63%"}}}, "EMEA": {"Direct Named": {"QSO": {"QTD": "1.2K", "Attain": "0.12%", "YoY": "1.28%"}, "Pipeline": {"QTD": "$3118.31M", "Attain": "2.61%", "YoY": "-0.74%"}}, "Partner": {"Pipeline": {"QTD": "$6256.96M", "Attain": "2.59%", "YoY": "0.25%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.97K", "Attain": "0.12%", "YoY": "1.1%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12485.45M", "Attain": "2.6%", "YoY": "0.02%"}}}, "LATAM": {"Direct Named": {"QSO": {"QTD": "1.16K", "Attain": "0.11%", "YoY": "-1.57%"}, "Pipeline": {"QTD": "$3118.28M", "Attain": "2.6%", "YoY": "0.05%"}}, "Partner": {"Pipeline": {"QTD": "$6152.72M", "Attain": "2.57%", "YoY": "-0.71%"}}, "GWS Direct QSOs": {"QSO": {"QTD": "1.95K", "Attain": "0.12%", "YoY": "-1.16%"}}, "GWS Direct + Partner Pipe": {"Pipeline": {"QTD": "$12434.49M", "Attain": "2.6%", "YoY": "-0.48%"}}}}, "insights": [{"title": "EMEA SMB Direct Named Program Underperforming in QTD QSO Pacing", "narrative": "The *EMEA SMB Direct Named* marketing program is significantly underperforming, achieving a QTD QSO pacing of only 0.17% against a target of 95%. This indicates a critical bottleneck in converting inquiries from key channels like *Display - Paid Social* (12.59% of inquiries) and *Email* (13.04% of inquiries) to QSOs.  A thorough analysis of campaign-level conversion rates within this program, particularly focusing on sales follow-up rates, is crucial to identify the root causes and implement corrective actions."}, {"title": "EMEA SMB Partner QTD QSO Pacing Lags Despite Slight YoY Improvement", "narrative": "EMEA SMB Partner QTD QSO Pacing is alarmingly low at **0.17%**, signaling potential difficulties in achieving quarterly targets despite a marginal **0.67%** YoY increase.  This underperformance is further emphasized by the substantial *$11.3M* pipeline generated by campaigns like *P&C Top Summit January 2024*, which unfortunately struggles to translate into qualified opportunities due to a low QSO conversion rate. To address this, prioritize optimizing pipeline conversion by analyzing high-performing campaigns like *'24 Gartner Supply Chain Symposium/Xpo'* (**50.63%** SAL Conversion Rate) and replicating their successful strategies within the EMEA SMB Partner segment. Additionally, benchmarking the performance of EMEA SMB Partner marketing programs against successful initiatives in other regions like NORTHAM or PUBLIC SECTOR, such as *Cloud Architecture Framework: Made in The Cloud*, can provide valuable insights for improvement."}]},
    23: {"data": {"NORTHAM": {"Ticket Volume": {"Q1": "83", "Chg": "65"}, "SLA Adherence": {"Q1": "97.59%", "Chg": "2.41%"}, "Marketer Satisfaction": {"Q1": "100.0%", "Chg": "0.0%"}}, "GLOBAL": {"Ticket Volume": {"Q1": "75", "Chg": "50"}, "SLA Adherence": {"Q1": "84.0%", "Chg": "16.0%"}, "Marketer Satisfaction": {"Q1": "100.0%", "Chg": "4.0%"}}, "LATAM": {"Ticket Volume": {"Q1": "30", "Chg": "25"}, "SLA Adherence": {"Q1": "76.67%", "Chg": "23.33%"}, "Marketer Satisfaction": {"Q1": "100.0%", "Chg": "0.0%"}}, "EMEA": {"Ticket Volume": {"Q1": "153", "Chg": "138"}, "SLA Adherence": {"Q1": "86.93%", "Chg": "13.07%"}, "Marketer Satisfaction": {"Q1": "100.0%", "Chg": "0.0%"}}, "APAC": {"Ticket Volume": {"Q1": "87", "Chg": "77"}, "SLA Adherence": {"Q1": "90.8%", "Chg": "9.2%"}, "Marketer Satisfaction": {"Q1": "98.85%", "Chg": "-1.15%"}}, "JAPAN": {"Ticket Volume": {"Q1": "13", "Chg": "7"}, "SLA Adherence": {"Q1": "92.31%", "Chg": "7.69%"}, "Marketer Satisfaction": {"Q1": "100.0%", "Chg": "0.0%"}}, "PUBLIC SECTOR": {"Ticket Volume": {"Q1": "37", "Chg": "36"}, "SLA Adherence": {"Q1": "86.49%", "Chg": "13.51%"}, "Marketer Satisfaction": {"Q1": "100.0%", "Chg": "0.0%"}}, "TOTAL": {"Ticket Volume": {"Q1": "522", "Chg": "430", "YTD": "855"}, "SLA Adherence": {"Q1": "88.89%", "Chg": "11.11%", "YTD": "90.53%"}, "Marketer Satisfaction": {"Q1": "99.62%", "Chg": "0.7%", "YTD": "99.42%"}}}, "insights": ["**AMP Issues:** Many tickets report problems with AMP, including submission errors, requests getting stuck, inability to create requests, and general unresponsiveness. This suggests potential bugs or usability issues within the AMP platform.", "**Knak Access:** Numerous users are requesting access to Knak, indicating a potential bottleneck in the user provisioning process or a lack of clarity regarding access criteria.", "**Dashboard & Reporting Issues:**  Multiple tickets highlight discrepancies in data across various dashboards (e.g., Demand Funnel, OKR tracker, GCM Health), impacting data analysis and decision-making. This suggests a need for data reconciliation and improved data integrity across platforms.", "**Lead Routing and Scoring:** Several issues pertain to leads not being routed or scored correctly, impacting campaign effectiveness and sales follow-up. This points to potential issues with lead routing rules, scoring models, or data quality within the CRM/marketing automation system."]},    
    38: {"data": {"Consideration": {"EMEA": {"value": "57.15%", "QoQ": "-"}, "TOTAL": {"value": "59.63%", "QoQ": "-"}, "JAPAC": {"value": "68.51%", "QoQ": "-"}, "LATAM": {"value": "84.21%", "QoQ": "-"}, "NORTHAM": {"value": "51.37%", "QoQ": "-"}}, "AI Perception": {"EMEA": {"value": "14.25%", "QoQ": "-"}, "TOTAL": {"value": "15.2%", "QoQ": "-"}, "JAPAC": {"value": "19.66%", "QoQ": "-"}, "LATAM": {"value": "26.86%", "QoQ": "-"}, "NORTHAM": {"value": "10.84%", "QoQ": "-"}}}, "insights": {"Consideration": [[{"title": "GCP's Lead Challenged by AWS in LATAM", "narrative": "Despite leading in LATAM with a KPI Value of 84.21%, GCP faces strong competition from AWS (91.4%) and Azure (83.5%), especially in the Financial Services and Healthcare & Life Sciences industries where AWS holds a higher KPI Value."}, {"title": "GCP's Performance in the German Automotive Industry", "narrative": "GCP's KPI Value in the German Automotive industry is lagging behind AWS, indicating a potential area for improvement and growth within this crucial sector."}, {"title": "GCP's Strong Showing in Education Sector", "narrative": "GCP demonstrates a strong KPI Value in the Education sector, particularly among institutions with 5,000-9,999 employees, highlighting its success in meeting the specific needs of this industry."}], {"title": "GCP's Strong JAPAC Presence", "narrative": "GCP demonstrates strength in JAPAC with a leading KPI Value of 68.51%, surpassing AWS (72.62%) and Azure (73.25%). This success is driven by strong performance in key industries such as [Mention specific industries with high KPI values in JAPAC], highlighting GCP's resonance with the region's unique market demands."}], "AI": [{"title": "GCP Dominates LATAM Cloud Market with 26.86% Mindshare", "narrative": "GCP is the leading cloud provider in LATAM with 26.86% mindshare, outperforming AWS by 10.02 percentage points and Azure by 5.19 percentage points. This strong performance is driven by GCP's dominance in key industries such as [mention specific industries with high GCP KPI values in LATAM]. For example, GCP holds [mention specific KPI percentage] share in the [specific industry] industry in LATAM."}, {"title": "GCP's Global Lag and LATAM Leadership", "narrative": "Despite a strong showing in LATAM, GCP lags behind Azure globally by 10.39 percentage points and trails AWS in all regions except LATAM. For instance, GCP only holds a 25% mindshare among Technology Execs in France within 1k-4.9k employee companies, while Azure dominates with 50%. This highlights a need to focus on key sectors and regions where GCP trails."}]}},
    39: {"data": {"Unaided Awareness": {"EMEA": {"value": "63.53%", "QoQ": "-"}, "TOTAL": {"value": "65.93%", "QoQ": "-"}, "JAPAC": {"value": "63.92%", "QoQ": "-"}, "LATAM": {"value": "72.11%", "QoQ": "-"}, "NORTHAM": {"value": "69.25%", "QoQ": "-"}}, "Familiarity": {"EMEA": {"value": "93.88%", "QoQ": "-"}, "TOTAL": {"value": "95.39%", "QoQ": "-"}, "JAPAC": {"value": "95.93%", "QoQ": "-"}, "LATAM": {"value": "97.01%", "QoQ": "-"}, "NORTHAM": {"value": "96.15%", "QoQ": "-"}}}, "insights": {"leading_indicators": [{"title": "GCP trails Azure significantly in EMEA", "narrative": "GCP trails Azure significantly in EMEA with 14.25% mindshare vs 26.79%, representing a 12.54 point difference. This is largely driven by France and Germany, representing (x% and y% respectively of EMEA responses).  **Key Takeaway:** GCP needs to increase brand visibility and consideration in EMEA, specifically within the French and German markets. **Next Steps:**  Consider allocating more marketing spend in EMEA to close the gap, and tailor messaging to address the specific needs and priorities of the French and German markets. "}, {"title": "GCP trails in Brand Familiarity", "narrative": "Despite a strong showing in LATAM, GCP lags behind Azure globally by 10.39 percentage points and trails AWS in all regions except LATAM. GCP\u2019s Familiarity is particularly low in EMEA and APAC, indicating a need for increased brand building efforts in these regions. **Key Takeaway:** GCP needs to prioritize brand building efforts in EMEA and APAC to close the familiarity gap with AWS and Azure. **Next Steps:**  Increase investments in brand awareness campaigns, thought leadership initiatives, and developer outreach programs in these regions."}]}}
}

@app.route("/")
def index():
    return "Lemur Service"

@app.route("/health", methods=["GET"])
def health_check():
    """
    Health check endpoint.
    ---
    responses:
      200:
        description: Service is healthy
    """
    return jsonify({"status": "healthy"}), 200

@app.route("/generate", methods=["POST"])
def generate():
    """
    Endpoint to generate a presentation.
    ---
    parameters:
      - name: data
        in: body
        required: true
        schema:
          type: object
          properties:
            file_id:
              type: string
              example: "1"
    responses:
      200:
        description: Presentation generated successfully
        schema:
          type: object
          properties:
            original_parameters:
              type: object
            presentation_link:
              type: string
            api_data:
              type: object
      500:
        description: Error generating presentation
    """
    try:
        data = request.get_json()
        logger.info(f"Received request data: {data}")

        slide_numbers = [14, 15, 16, 17, 23, 38, 39]
        slide_data = {}
        api_data = {}

        # api_url = "http://34.90.192.243/insight_slide"
        # for slide_no in slide_numbers:
        #     slide_data[slide_no] = fetch_slide_data_with_retry(api_url, slide_no)
        #     api_data[slide_no] = slide_data[slide_no]
        #     logger.info(
        #         f"Received slide data from API for slide {slide_no}: {slide_data[slide_no]}"
        #     )

        # Use hardcoded data for each slide number
        for slide_no in slide_numbers:
            slide_data[slide_no] = hardcoded_data.get(slide_no, {})
            api_data[slide_no] = slide_data[slide_no]
            logger.info(
                f"Using hardcoded data for slide {slide_no}: {slide_data[slide_no]}"
            )

        # Generate the presentation
        presentation_link = create_presentation(slide_data, data["file_id"])
        logger.info(f"Generated presentation link: {presentation_link}")

        response_data = {
            "original_parameters": data,
            "presentation_link": presentation_link,
            "api_data": api_data,  # Add API data to response
        }
        return jsonify(response_data), 200
    except Exception as e:
        logger.error(f"Error generating presentation: {e}")
        return jsonify({"error": str(e)}), 500

def create_presentation(data, file_id):
    """
    Create a presentation and populate it with data.
    """
    try:
        # Load the template presentation
        template_path = "template.pptx"
        prs = Presentation(template_path)

        # Populate the presentation with data
        for slide_no, content in data.items():
            if slide_no - 1 < len(prs.slides):
                slide = prs.slides[
                    slide_no - 1
                ]  # Adjust index since slides are 0-indexed
                populate_slide(slide, content, slide_no)
            else:
                logger.error(
                    f"Slide number {slide_no} is out of range for the presentation"
                )

        # Save the modified presentation
        output_path = f"/tmp/{file_id}.pptx"
        prs.save(output_path)

        # Upload the presentation to Google Drive
        file_metadata = {"name": f"Generated Presentation {file_id}"}
        media = MediaFileUpload(
            output_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        uploaded_file = upload_to_drive_with_retry(file_metadata, media)
        logger.info(f"Uploaded presentation with ID: {uploaded_file.get('id')}")

        # Set file permissions to make it accessible by anyone with the link
        permission = {
            "type": "anyone",
            "role": "reader",
        }
        drive_service.permissions().create(
            fileId=uploaded_file["id"],
            body=permission,
        ).execute()

        return f"https://drive.google.com/file/d/{uploaded_file.get('id')}/view"
    except Exception as e:
        logger.error(f"Error creating presentation: {e}")
        raise e  # Re-raise the exception after logging it

def upload_to_drive_with_retry(file_metadata, media, retries=3):
    attempt = 0
    while attempt < retries:
        try:
            uploaded_file = (
                drive_service.files()
                .create(body=file_metadata, media_body=media, fields="id")
                .execute()
            )
            return uploaded_file
        except Exception as e:
            logger.error(f"Attempt {attempt + 1} failed with error: {e}")
            attempt += 1
            time.sleep(2**attempt)  # Exponential backoff
    raise Exception("Failed to upload file after several retries")

def set_font(cell, font_name="Arial", font_size=8):
    """
    Set the font for a table cell.
    """
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)

def set_yoy_color(cell, yoy_value):
    """
    Set the color for the YoY cell based on its value.
    """
    try:
        if yoy_value:
            yoy = float(yoy_value.strip("%"))
            if yoy > 100:
                color = RGBColor(0, 255, 0)  # Green
            elif 90 <= yoy <= 100:
                color = RGBColor(218, 165, 32)  # Yellow goldenrod
            else:
                color = RGBColor(255, 0, 0)  # Red
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = color
        else:
            logger.error("Empty YoY value provided.")
    except ValueError:
        logger.error(f"Invalid YoY value: {yoy_value}")

def populate_slide(slide, content, slide_number):
    """
    Populate a slide with the given content.
    """
    try:
        main_table = None
        insights_table = None
        table_count = 0

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                if len(table.rows) <= 4:
                    insights_table = table
                else:
                    main_table = table
                table_count += 1

        if main_table:
            logger.info(
                f"Main table dimensions: {len(main_table.rows)} rows x {len(main_table.columns)} columns"
            )

            if slide_number in [14, 15, 16, 17]:
                regions = [
                    "NORTHAM",
                    "LATAM",
                    "EMEA",
                    "JAPAC",
                    "US PUBLIC SECTOR",
                    "GLOBAL",
                ]

                if slide_number == 14:
                    metrics = [
                        ("Direct Named", "QSO", "QTD", "Attain", "YoY"),
                        ("Direct Named", "Pipeline", "QTD", "Attain", "YoY"),
                        ("Startup", "QSO", "QTD", "Attain", "YoY"),
                        ("Startup", "Pipeline", "QTD", "Attain", "YoY"),
                        ("SMB", "QSO", "QTD", "Attain", "YoY"),
                        ("SMB", "Pipeline", "QTD", "Attain", "YoY"),
                        ("Partner", "Pipeline", "QTD", "Attain", "YoY"),
                        ("GCP Direct QSOs", "QSO", "QTD", "Attain", "YoY"),
                        ("GCP Direct + Partner Pipe", "Pipeline", "QTD", "Attain", "YoY"),
                    ]
                elif slide_number == 15:
                    metrics = [
                        ("Direct Named", "QSO", "QTD", "Attain", "YoY"),
                        ("Direct Named", "Pipeline", "QTD", "Attain", "YoY"),
                        ("Startup", "QSO", "QTD", "Attain", "YoY"),
                        ("Startup", "Pipeline", "QTD", "Attain", "YoY"),
                        ("SMB", "QSO", "QTD", "Attain", "YoY"),
                        ("SMB", "Pipeline", "QTD", "Attain", "YoY"),
                        ("Partner", "Pipeline", "QTD", "Attain", "YoY"),
                        ("GCP Direct QSOs", "QSO", "QTD", "Attain", "YoY"),
                        ("GCP Direct + Partner Pipe", "Pipeline", "QTD", "Attain", "YoY"),
                    ]
                elif slide_number in [16, 17]:
                    metrics = [
                        ("Direct Named", "QSO", "QTD", "Attain", "YoY"),
                        ("Direct Named", "Pipeline", "QTD", "Attain", "YoY"),
                        ("Partner", "Pipeline", "QTD", "Attain", "YoY"),
                        ("GWS Direct QSOs", "QSO", "QTD", "Attain", "YoY"),
                        ("GWS Direct + Partner Pipe", "Pipeline", "QTD", "Attain", "YoY"),
                    ]

                start_row = 3

                for i, (metric_category, metric_type, qtd_key, attain_key, yoy_key) in enumerate(metrics):
                    for j, region in enumerate(regions):
                        metric_data = (
                            content.get("data", {})
                            .get(region, {})
                            .get(metric_category, {})
                            .get(metric_type, {})
                        )
                        region_value_qtd_attain = (
                            f"{metric_data.get(qtd_key, '')} ({metric_data.get(attain_key, '')})"
                            if metric_data.get(qtd_key) or metric_data.get(attain_key)
                            else ""
                        )
                        region_value_yoy = metric_data.get(yoy_key, "")

                        # QTD and Attain
                        cell = main_table.cell(start_row + i, 1 + (j * 2))
                        cell.text = region_value_qtd_attain
                        set_font(cell)

                        # YoY
                        cell = main_table.cell(start_row + i, 2 + (j * 2))
                        cell.text = region_value_yoy
                        set_font(cell)
                        # set_yoy_color(cell, region_value_yoy)

                        logger.info(
                            f"Populated {metric_category} {metric_type} for {region} with QTD+Attain: {region_value_qtd_attain} and YoY: {region_value_yoy}"
                        )

            elif slide_number == 23:
                regions = ["NORTHAM", "GLOBAL", "LATAM", "EMEA", "APAC", "JAPAN", "PUBLIC SECTOR", "TOTAL"]
                metrics = [
                    ("Ticket Volume", "Q1", "Chg", "YTD"),
                    ("SLA Adherence", "Q1", "Chg", "YTD"),
                    ("Marketer Satisfaction", "Q1", "Chg", "YTD")
                ]

                start_row = 2

                for i, (metric, q1_key, chg_key, ytd_key) in enumerate(metrics):
                    for j, region in enumerate(regions):
                        region_data = content.get("data", {}).get(region, {}).get(metric, {})
                        q1_value = region_data.get(q1_key, "")
                        chg_value = region_data.get(chg_key, "")
                        ytd_value = region_data.get(ytd_key, "") if region == "TOTAL" else ""

                        # Q1
                        cell = main_table.cell(start_row + i, 1 + (j * 2))
                        cell.text = q1_value
                        set_font(cell)

                        # Chg
                        cell = main_table.cell(start_row + i, 2 + (j * 2))
                        cell.text = chg_value
                        set_font(cell)

                        # YTD (only for TOTAL column)
                        if region == "TOTAL":
                            cell = main_table.cell(start_row + i, 3 + (j * 2))
                            cell.text = ytd_value
                            set_font(cell)

                        logger.info(
                            f"Populated {metric} for {region} with Q1: {q1_value}, Chg: {chg_value}, and YTD: {ytd_value}"
                        )

            elif slide_number == 38:
                categories = ["Consideration", "AI Perception"]
                regions = ["EMEA", "TOTAL", "JAPAC", "LATAM", "NORTHAM"]

                start_row = 3

                for i, category in enumerate(categories):
                    for j, region in enumerate(regions):
                        region_data = content.get("data", {}).get(category, {}).get(region, {})
                        value = region_data.get("value", "")
                        qoq = region_data.get("QoQ", "")

                        # Value
                        cell = main_table.cell(start_row + i, 1 + (j * 2))
                        cell.text = value
                        set_font(cell)

                        # QoQ
                        cell = main_table.cell(start_row + i, 2 + (j * 2))
                        cell.text = qoq
                        set_font(cell)

                        logger.info(
                            f"Populated {category} for {region} with Value: {value} and QoQ: {qoq}"
                        )

            elif slide_number == 39:
                categories = ["Unaided Awareness", "Familiarity"]
                regions = ["EMEA", "TOTAL", "JAPAC", "LATAM", "NORTHAM"]

                start_row = 3

                for i, category in enumerate(categories):
                    for j, region in enumerate(regions):
                        region_data = content.get("data", {}).get(category, {}).get(region, {})
                        value = region_data.get("value", "")
                        qoq = region_data.get("QoQ", "")

                        # Value
                        cell = main_table.cell(start_row + i, 1 + (j * 2))
                        cell.text = value
                        set_font(cell)

                        # QoQ
                        cell = main_table.cell(start_row + i, 2 + (j * 2))
                        cell.text = qoq
                        set_font(cell)

                        logger.info(
                            f"Populated {category} for {region} with Value: {value} and QoQ: {qoq}"
                        )

        if insights_table:
            logger.info(
                f"Insights table dimensions: {len(insights_table.rows)} rows x {len(insights_table.columns)} columns"
            )

            # Populate insights table differently based on slide number
            if slide_number in [14, 15, 16, 17, 23]:
                drivers = content.get("drivers", [])
                insights = content.get("insights", [])
                logger.info(f"Length of insights is: {len(insights)}")
                logger.info(f"Length of insight table is: {len(insights_table.rows)}")

                for i in range(
                    max(len(drivers), len(insights))
                ):  # Ensure we loop through the longest list
                    if i < len(insights_table.rows):
                        driver_text = drivers[i] if i < len(drivers) else ""
                        insight_text = insights[i] if i < len(insights) else ""
                        if isinstance(insight_text, str):
                            try:
                                insight_text = json.loads(insight_text)
                            except json.JSONDecodeError as e:
                                logger.error(f"Error parsing insight string: {e}")
                                insight_text = {}
                        driver_text = insight_text.get("title", "")
                        insight_text = insight_text.get("narrative", "")

                        cell = insights_table.cell(i, 0)
                        set_font(cell)

                        # Bold the driver_text part
                        paragraph = cell.text_frame.paragraphs[0]
                        run = paragraph.add_run()
                        run.text = driver_text
                        run.font.bold = True
                        run.font.size = Pt(8)
                        run.font.name = "Arial"
                        run = paragraph.add_run()
                        run.text = f" {insight_text}"
                        run.font.size = Pt(8)
                        run.font.name = "Arial"
                        logger.info(
                            f"Populated insight {i + 1}: {driver_text} {insight_text}"
                        )

                recommendations = content.get("recommendations", [])
                recommendations = [re.sub(r"\*\*", "", rec) for rec in recommendations]
                footnote_text = " ".join(recommendations)
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text += "\n" + footnote_text

            elif slide_number == 38:
                insights = content.get("insights", {})
                logger.info(f"Insights for slide 38: {insights}")

                for category, category_insights in insights.items():
                    if isinstance(category_insights, list):  # Ensure it's a list
                        for i, insight_group in enumerate(category_insights):
                            if isinstance(insight_group, list):  # Ensure it's a list
                                for j, insight in enumerate(insight_group):
                                    if isinstance(insight, dict):  # Ensure it's a dict
                                        if i * len(insight_group) + j < len(insights_table.rows):
                                            cell = insights_table.cell(i * len(insight_group) + j, 0)
                                            set_font(cell)

                                            title = insight.get("title", "")
                                            narrative = insight.get("narrative", "")

                                            # Bold the title part
                                            paragraph = cell.text_frame.paragraphs[0]
                                            run = paragraph.add_run()
                                            run.text = title
                                            run.font.bold = True
                                            run.font.size = Pt(8)
                                            run.font.name = "Arial"
                                            run = paragraph.add_run()
                                            run.text = f" {narrative}"
                                            run.font.size = Pt(8)
                                            run.font.name = "Arial"
                                            logger.info(
                                                f"Populated insight {i + 1}: {title} {narrative}"
                                            )

                
            elif slide_number == 39:
                insights = content.get("insights", {}).get("leading_indicators", [])
                logger.info(f"Leading indicators for slide 39: {insights}")

                for i, insight in enumerate(insights):
                    if isinstance(insight, dict):  # Ensure it's a dict
                        if i < len(insights_table.rows):
                            cell = insights_table.cell(i, 0)
                            set_font(cell)

                            title = insight.get("title", "")
                            narrative = insight.get("narrative", "")

                            # Bold the title part
                            paragraph = cell.text_frame.paragraphs[0]
                            run = paragraph.add_run()
                            run.text = title
                            run.font.bold = True
                            run.font.size = Pt(8)
                            run.font.name = "Arial"
                            run = paragraph.add_run()
                            run.text = f" {narrative}"
                            run.font.size = Pt(8)
                            run.font.name = "Arial"
                            logger.info(
                                f"Populated insight {i + 1}: {title} {narrative}"
                            )


    except Exception as e:
        logger.error(f"Error populating slide: {e}")
        raise e  # Re-raise the exception after logging it


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8080)
